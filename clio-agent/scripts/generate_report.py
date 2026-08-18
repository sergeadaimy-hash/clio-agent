#!/usr/bin/env python3
"""CLIO daily report generator: reads brand & report config from DB.

Usage: python3 scripts/generate_report.py YYYY-MM-DD
Prints the output PPTX path on stdout on success.
"""

import os
import sys
import json
import sqlite3
from io import BytesIO
from datetime import datetime

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ── Paths ───────────────────────────────────────────────────
ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
DB_PATH = os.path.join(ROOT, "db", "clio.db")
TEMPLATE_PATH = os.path.join(ROOT, "templates", "base_template.pptx")
REPORTS_ROOT = os.path.join(ROOT, "reports")

# ── Status colors (semantic, always fixed) ─────────────────
GREEN = RGBColor(0x4A, 0xDE, 0x80)
AMBER = RGBColor(0xFB, 0x92, 0x3C)
RED = RGBColor(0xF8, 0x71, 0x71)


def hex_to_rgb(hx):
    hx = (hx or "#3B82F6").lstrip("#")
    if len(hx) != 6:
        return RGBColor(0x3B, 0x82, 0xF6)
    return RGBColor(int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))


def rgb_to_hex(c):
    return f"#{c[0]:02X}{c[1]:02X}{c[2]:02X}" if hasattr(c, '__getitem__') else "#0F172A"


# ── Config from DB ──────────────────────────────────────────
def load_config():
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    brand = {}
    report = {}
    for key in ('brand_config', 'report_config', 'event_name', 'event_edition'):
        cur.execute("SELECT value FROM settings WHERE key = ?", (key,))
        row = cur.fetchone()
        if row and row[0]:
            if key == 'brand_config':
                try: brand = json.loads(row[0])
                except: pass
            elif key == 'report_config':
                try: report = json.loads(row[0])
                except: pass
            elif key == 'event_name':
                brand['event_name'] = row[0]
            elif key == 'event_edition':
                brand['event_edition'] = row[0]
    conn.close()
    return brand, report


class Colors:
    def __init__(self, brand):
        self.BG = hex_to_rgb(brand.get('background_color', '#0F172A'))
        self.TEXT = hex_to_rgb(brand.get('text_color', '#F1F5F9'))
        self.MUTED = hex_to_rgb(brand.get('muted_color', '#94A3B8'))
        self.PANEL = hex_to_rgb(brand.get('panel_color', '#1E293B'))
        self.ACCENT = hex_to_rgb(brand.get('primary_color', '#3B82F6'))
        self.bg_hex = brand.get('background_color', '#0F172A')
        self.accent_hex = brand.get('primary_color', '#3B82F6')
        self.panel_hex = brand.get('panel_color', '#1E293B')


# ── DB ──────────────────────────────────────────────────────
def fetch_data(date):
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cur = conn.cursor()
    cur.execute("SELECT * FROM departments ORDER BY id")
    departments = [dict(r) for r in cur.fetchall()]
    cur.execute(
        """SELECT ds.*, d.name AS department_name, d.stream_color, d.head_name
           FROM daily_submissions ds
           JOIN departments d ON d.id = ds.department_id
           WHERE ds.submission_date = ? AND ds.is_submitted = 1
           ORDER BY ds.submitted_at""",
        (date,),
    )
    submissions = []
    for r in cur.fetchall():
        d = dict(r)
        try: d["schedule_updates"] = json.loads(d.get("schedule_updates") or "[]")
        except: d["schedule_updates"] = []
        try: d["photos"] = json.loads(d.get("photos") or "[]")
        except: d["photos"] = []
        try: d["photo_captions"] = json.loads(d.get("photo_captions") or "{}")
        except: d["photo_captions"] = {}
        submissions.append(d)
    conn.close()
    return departments, submissions


# ── Slide helpers ───────────────────────────────────────────
def new_presentation(report):
    if os.path.exists(TEMPLATE_PATH):
        try:
            return Presentation(TEMPLATE_PATH)
        except Exception as e:
            sys.stderr.write(f"template load failed, using blank: {e}\n")
    prs = Presentation()
    dims = report.get('slide_dimensions', '16:9')
    if dims == '4:3':
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs, C):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = C.BG
    return slide


def add_text(slide, left, top, width, height, text, *, size=14, bold=False, color=None, align=PP_ALIGN.LEFT, font_name="Calibri"):
    if color is None:
        color = RGBColor(0xF1, 0xF5, 0xF9)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text if text is not None else ""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_rect(slide, left, top, width, height, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        shp.line.fill.background()
    return shp


def add_slide_number(slide, idx, total, C, font):
    add_text(slide, Inches(12.6), Inches(7.1), Inches(0.7), Inches(0.3),
             f"{idx} / {total}", size=9, color=C.MUTED, align=PP_ALIGN.RIGHT, font_name=font)


def add_footer_accent(slide, color):
    add_rect(slide, 0, 0, Inches(0.1), Inches(7.5), color)


def add_logo(slide, brand):
    logo_path = brand.get('logo_path', '')
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, Inches(11.5), Inches(0.25), width=Inches(1.5))
        except Exception as e:
            sys.stderr.write(f"logo err: {e}\n")


# ── Slide 1: overview ───────────────────────────────────────
def slide_overview(prs, date, departments, submissions, brand, C, report, font):
    slide = blank_slide(prs, C)
    add_logo(slide, brand)

    event_name = brand.get('event_name', os.environ.get("EVENT_NAME", "Event"))
    edition = brand.get('event_edition', os.environ.get("EVENT_EDITION", ""))
    company = brand.get('company_name', '')

    # Company name top-left
    if company:
        add_text(slide, Inches(0.4), Inches(0.2), Inches(4), Inches(0.3),
                 company.upper(), size=10, bold=True, color=C.ACCENT, font_name=font)

    header_y = Inches(0.5) if company else Inches(0.4)
    add_text(slide, Inches(0.4), header_y, Inches(10), Inches(0.7),
             f"{event_name}  •  {edition}".strip(" •"), size=28, bold=True, color=C.TEXT, font_name=font)
    add_text(slide, Inches(0.4), header_y + Inches(0.7), Inches(10), Inches(0.4),
             f"Daily Progress Report  •  Generated by CLIO  •  {date}", size=14, color=C.MUTED, font_name=font)

    rcfg = report.get('overview', {})
    show_badges = rcfg.get('show_badges', True)
    show_average = rcfg.get('show_average', True)

    grid_top = Inches(1.8)
    col_w = Inches(6.2)
    row_h = Inches(0.75)
    gap = Inches(0.1)
    for i, dept in enumerate(departments):
        col = i % 2
        row = i // 2
        left = Inches(0.4) + (col_w + gap) * col
        top = grid_top + (row_h + gap) * row

        add_rect(slide, left, top, col_w, row_h, C.PANEL)
        add_rect(slide, left, top, Inches(0.12), row_h, hex_to_rgb(dept["stream_color"]))
        add_text(slide, left + Inches(0.25), top + Inches(0.08),
                 Inches(3.0), Inches(0.35), dept["name"], size=14, bold=True,
                 color=hex_to_rgb(dept["stream_color"]), font_name=font)

        sub = next((s for s in submissions if s["department_id"] == dept["id"]), None)
        if sub:
            pct = int(sub.get("overall_progress") or 0)
            add_text(slide, left + Inches(3.4), top + Inches(0.1),
                     Inches(1.4), Inches(0.5), f"{pct}%", size=22, bold=True, color=C.TEXT, font_name=font)
            if show_badges:
                if pct >= 80:
                    badge_color, badge_text = GREEN, "ON TRACK"
                elif pct >= 50:
                    badge_color, badge_text = AMBER, "IN PROGRESS"
                else:
                    badge_color, badge_text = RED, "AT RISK"
                add_text(slide, left + Inches(4.8), top + Inches(0.2),
                         Inches(1.3), Inches(0.35), badge_text, size=10, bold=True, color=badge_color, font_name=font)
        else:
            add_text(slide, left + Inches(3.4), top + Inches(0.2),
                     Inches(2.7), Inches(0.35), "No Data Submitted", size=11, color=C.MUTED, font_name=font)

    if show_average:
        if submissions:
            avg = round(sum((s.get("overall_progress") or 0) for s in submissions) / len(submissions))
        else:
            avg = 0
        add_text(slide, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.5),
                 f"Overall Average Completion: {avg}%  •  {len(submissions)}/{len(departments)} departments reporting",
                 size=14, color=C.MUTED, font_name=font)
    return slide


# ── Department overview ────────────────────────────────────
def slide_dept_overview(prs, dept, sub, brand, C, report, font):
    slide = blank_slide(prs, C)
    color = hex_to_rgb(dept["stream_color"])
    add_footer_accent(slide, color)
    add_logo(slide, brand)

    rcfg = report.get('dept_overview', {})

    add_text(slide, Inches(0.4), Inches(0.4), Inches(10), Inches(0.5),
             dept["name"].upper(), size=26, bold=True, color=color, font_name=font)
    add_text(slide, Inches(0.4), Inches(0.95), Inches(10), Inches(0.4),
             f"HOD: {dept.get('head_name') or 'n/a'}  •  Submitted: {sub.get('submitted_at','')}",
             size=12, color=C.MUTED, font_name=font)

    pct = int(sub.get("overall_progress") or 0)
    add_text(slide, Inches(0.4), Inches(1.6), Inches(6), Inches(1.4),
             f"{pct}%", size=96, bold=True, color=C.TEXT, font_name=font)
    add_text(slide, Inches(0.4), Inches(3.0), Inches(6), Inches(0.4),
             "Overall Completion", size=12, color=C.MUTED, font_name=font)

    add_text(slide, Inches(6.7), Inches(1.6), Inches(6.3), Inches(0.3),
             "STATUS", size=10, bold=True, color=C.MUTED, font_name=font)
    add_text(slide, Inches(6.7), Inches(1.9), Inches(6.3), Inches(2.0),
             sub.get("status_text") or "n/a", size=13, color=C.TEXT, font_name=font)

    if rcfg.get('show_highlights', True) and sub.get("highlights"):
        add_rect(slide, Inches(0.4), Inches(3.6), Inches(0.08), Inches(1.7), GREEN)
        add_text(slide, Inches(0.6), Inches(3.6), Inches(5.8), Inches(0.3),
                 "HIGHLIGHTS", size=10, bold=True, color=GREEN, font_name=font)
        add_text(slide, Inches(0.6), Inches(3.9), Inches(5.8), Inches(1.5),
                 sub["highlights"], size=12, color=C.TEXT, font_name=font)

    if rcfg.get('show_blockers', True) and sub.get("blockers"):
        add_rect(slide, Inches(6.7), Inches(3.6), Inches(0.08), Inches(1.7), RED)
        add_text(slide, Inches(6.9), Inches(3.6), Inches(6.1), Inches(0.3),
                 "BLOCKERS", size=10, bold=True, color=RED, font_name=font)
        add_text(slide, Inches(6.9), Inches(3.9), Inches(6.1), Inches(1.5),
                 sub["blockers"], size=12, color=C.TEXT, font_name=font)

    return slide


# ── Schedule + donut ───────────────────────────────────────
def make_donut(schedule, color_hex, bg_hex, text_hex, muted_hex):
    completed = sum(1 for a in schedule if (a.get("status") or "").lower() == "completed")
    in_prog = sum(1 for a in schedule if (a.get("status") or "").lower() == "in progress")
    pending = sum(1 for a in schedule if (a.get("status") or "").lower() == "pending")
    total = completed + in_prog + pending
    if total == 0:
        values, labels, colors = [1], ["No Data"], ["#374151"]
    else:
        values = [completed, in_prog, pending]
        labels = [f"Completed ({completed})", f"In Progress ({in_prog})", f"Pending ({pending})"]
        colors = [color_hex, "#FB923C", "#374151"]

    fig, ax = plt.subplots(figsize=(4.2, 3.2), dpi=150)
    fig.patch.set_facecolor(bg_hex)
    ax.set_facecolor(bg_hex)
    wedges, _ = ax.pie(values, colors=colors, startangle=90,
                       wedgeprops=dict(width=0.38, edgecolor=bg_hex))
    pct_done = round((completed / total) * 100) if total else 0
    ax.text(0, 0.05, f"{pct_done}%", ha="center", va="center",
            color=text_hex, fontsize=26, fontweight="bold")
    ax.text(0, -0.22, "done", ha="center", va="center", color=muted_hex, fontsize=10)
    ax.legend(labels, loc="lower center", bbox_to_anchor=(0.5, -0.1),
              frameon=False, labelcolor=text_hex, fontsize=9, ncol=1)
    buf = BytesIO()
    plt.savefig(buf, format="png", facecolor=fig.get_facecolor(), bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def slide_dept_schedule(prs, dept, sub, brand, C, report, font):
    slide = blank_slide(prs, C)
    color = hex_to_rgb(dept["stream_color"])
    add_footer_accent(slide, color)
    add_logo(slide, brand)

    rcfg = report.get('dept_schedule', {})
    show_donut = rcfg.get('show_donut', True)

    add_text(slide, Inches(0.4), Inches(0.4), Inches(10), Inches(0.5),
             f"{dept['name'].upper()}  •  TODAY'S SCHEDULE", size=22, bold=True, color=color, font_name=font)

    schedule = sub.get("schedule_updates") or []
    table_width = Inches(8.0) if show_donut else Inches(12.5)

    header_top = Inches(1.3)
    header_h = Inches(0.4)
    add_rect(slide, Inches(0.4), header_top, table_width, header_h, C.PANEL)
    add_text(slide, Inches(0.5), header_top + Inches(0.07), Inches(1.2), Inches(0.3),
             "TIME", size=10, bold=True, color=C.MUTED, font_name=font)
    add_text(slide, Inches(1.8), header_top + Inches(0.07), Inches(4.5), Inches(0.3),
             "ACTIVITY", size=10, bold=True, color=C.MUTED, font_name=font)
    add_text(slide, Inches(6.4), header_top + Inches(0.07), Inches(1.9), Inches(0.3),
             "STATUS", size=10, bold=True, color=C.MUTED, font_name=font)

    row_h = Inches(0.38)
    alt_bg = hex_to_rgb(C.bg_hex.replace('#', ''))  # slightly different shade
    for i, act in enumerate(schedule[:12]):
        top = header_top + header_h + row_h * i
        if i % 2 == 0:
            add_rect(slide, Inches(0.4), top, table_width, row_h, RGBColor(0x0A, 0x0C, 0x12))
        add_text(slide, Inches(0.5), top + Inches(0.07), Inches(1.2), Inches(0.3),
                 act.get("time", ""), size=11, color=C.TEXT, font_name=font)
        add_text(slide, Inches(1.8), top + Inches(0.07), Inches(4.5), Inches(0.3),
                 (act.get("activity", "") or "")[:80], size=11, color=C.TEXT, font_name=font)
        st = (act.get("status") or "").strip()
        st_color = GREEN if st.lower() == "completed" else (AMBER if st.lower() == "in progress" else C.MUTED)
        add_text(slide, Inches(6.4), top + Inches(0.07), Inches(1.9), Inches(0.3),
                 st or "n/a", size=11, bold=True, color=st_color, font_name=font)

    if show_donut:
        try:
            chart_buf = make_donut(schedule, dept["stream_color"], C.bg_hex,
                                   brand.get('text_color', '#F1F5F9'),
                                   brand.get('muted_color', '#94A3B8'))
            slide.shapes.add_picture(chart_buf, Inches(8.8), Inches(1.3), width=Inches(4.2), height=Inches(4.2))
        except Exception as e:
            sys.stderr.write(f"chart err: {e}\n")

    return slide


# ── Photo wall ─────────────────────────────────────────────
def slide_photo_wall(prs, dept, photos, page_idx, page_total, brand, C, report, font, captions=None):
    slide = blank_slide(prs, C)
    color = hex_to_rgb(dept["stream_color"])
    add_footer_accent(slide, color)
    add_logo(slide, brand)

    rcfg = report.get('dept_photos', {})
    show_timestamps = rcfg.get('show_timestamps', True)

    title = f"{dept['name'].upper()}  •  PROGRESS PHOTOS"
    if page_total > 1:
        title += f"  ({page_idx}/{page_total})"
    add_text(slide, Inches(0.4), Inches(0.4), Inches(12), Inches(0.5),
             title, size=22, bold=True, color=color, font_name=font)

    per_page = len(photos)
    if per_page <= 4:
        cols, rows = 2, 2
    elif per_page <= 6:
        cols, rows = 3, 2
    else:
        cols, rows = 3, 3

    total_w = Inches(12.5)
    total_h = Inches(5.5)
    gap = Inches(0.15)
    cell_w = (total_w - gap * (cols - 1)) // cols
    cell_h = (total_h - gap * (rows - 1)) // rows
    x0 = Inches(0.4)
    y0 = Inches(1.3)

    for idx, photo_path in enumerate(photos):
        r = idx // cols
        c = idx % cols
        left = x0 + (cell_w + gap) * c
        top = y0 + (cell_h + gap) * r
        if not os.path.exists(photo_path):
            add_text(slide, left, top, cell_w, cell_h, "(missing)", size=10, color=C.MUTED, font_name=font)
            continue
        try:
            with Image.open(photo_path) as im:
                iw, ih = im.size
            ar = iw / ih if ih else 1.0
            cell_ar = cell_w / cell_h
            if ar > cell_ar:
                w = cell_w
                h = int(cell_w / ar)
            else:
                h = cell_h
                w = int(cell_h * ar)
            px = left + (cell_w - w) // 2
            py = top + (cell_h - h) // 2
            slide.shapes.add_picture(photo_path, px, py, width=w, height=h)
        except Exception as e:
            sys.stderr.write(f"photo err {photo_path}: {e}\n")
            continue

        # Caption and/or timestamp under photo
        caption_parts = []
        if captions:
            # Match by filename (captions keyed by original filename)
            base = os.path.basename(photo_path)
            for cap_key, cap_val in captions.items():
                if cap_key in base or base.endswith(cap_key):
                    caption_parts.append(cap_val)
                    break
        if show_timestamps:
            base = os.path.basename(photo_path)
            try:
                parts = base.split("_")
                for p in parts:
                    if len(p) == 8 and p[2] == "-" and p[5] == "-":
                        caption_parts.append(f"{p[:2]}:{p[3:5]}")
                        break
            except Exception:
                pass
        if caption_parts:
            caption_text = "  ·  ".join(caption_parts) if len(caption_parts) > 1 else caption_parts[0]
            add_text(slide, left, top + cell_h + Inches(0.02),
                     cell_w, Inches(0.35), caption_text, size=9, color=C.MUTED, font_name=font)

    return slide


def slide_no_submission(prs, dept, C, font):
    slide = blank_slide(prs, C)
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, RGBColor(0x12, 0x14, 0x1A))
    add_text(slide, Inches(0.4), Inches(3.0), Inches(12.5), Inches(1.0),
             f"{dept['name']}: No Update Submitted",
             size=32, bold=True, color=C.MUTED, align=PP_ALIGN.CENTER, font_name=font)
    add_text(slide, Inches(0.4), Inches(4.0), Inches(12.5), Inches(0.5),
             "HOD was reminded at 9:00 PM. No data received before report generation.",
             size=14, color=C.MUTED, align=PP_ALIGN.CENTER, font_name=font)
    return slide


# ── Main build ──────────────────────────────────────────────
def build(date):
    brand, report = load_config()
    C = Colors(brand)
    font = brand.get('font_name', 'Calibri') or 'Calibri'

    departments, submissions = fetch_data(date)
    prs = new_presentation(report)

    submitted_dept_ids = [s["department_id"] for s in submissions]
    submitted_depts_ordered = []
    for s in submissions:
        dept = next((d for d in departments if d["id"] == s["department_id"]), None)
        if dept:
            submitted_depts_ordered.append((dept, s))
    not_submitted = [d for d in departments if d["id"] not in submitted_dept_ids]

    slide_order = report.get('slide_order', ['overview', 'dept_overview', 'dept_schedule', 'dept_photos', 'no_submission'])

    # Overview
    if 'overview' in slide_order and report.get('overview', {}).get('enabled', True):
        slide_overview(prs, date, departments, submissions, brand, C, report, font)

    # Per-department slides
    for dept, sub in submitted_depts_ordered:
        for slide_type in slide_order:
            if slide_type == 'dept_overview' and report.get('dept_overview', {}).get('enabled', True):
                slide_dept_overview(prs, dept, sub, brand, C, report, font)
            elif slide_type == 'dept_schedule' and report.get('dept_schedule', {}).get('enabled', True):
                slide_dept_schedule(prs, dept, sub, brand, C, report, font)
            elif slide_type == 'dept_photos' and report.get('dept_photos', {}).get('enabled', True):
                photos = sub.get("photos") or []
                if photos:
                    per_page = report.get('dept_photos', {}).get('photos_per_page', 6)
                    pages = (len(photos) + per_page - 1) // per_page
                    for i in range(pages):
                        chunk = photos[i * per_page:(i + 1) * per_page]
                        captions = sub.get("photo_captions") or {}
                        slide_photo_wall(prs, dept, chunk, i + 1, pages, brand, C, report, font, captions)

    # No-submission notices
    if 'no_submission' in slide_order and report.get('no_submission', {}).get('enabled', True):
        for dept in not_submitted:
            slide_no_submission(prs, dept, C, font)

    # Number all slides
    total = len(prs.slides)
    for idx, slide in enumerate(prs.slides, start=1):
        add_slide_number(slide, idx, total, C, font)

    # Save
    out_dir = os.path.join(REPORTS_ROOT, date)
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "daily_report.pptx")
    prs.save(out_path)
    return out_path


def main():
    if len(sys.argv) < 2:
        sys.stderr.write("usage: generate_report.py YYYY-MM-DD\n")
        sys.exit(1)
    date = sys.argv[1]
    try:
        datetime.strptime(date, "%Y-%m-%d")
    except ValueError:
        sys.stderr.write(f"invalid date: {date}\n")
        sys.exit(1)
    try:
        out = build(date)
        print(out)
    except Exception as e:
        sys.stderr.write(f"generation failed: {e}\n")
        import traceback
        traceback.print_exc(file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
