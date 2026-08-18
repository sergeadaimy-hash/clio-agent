#!/usr/bin/env python3
"""CLIO weekly report generator.

Usage: python3 scripts/generate_weekly_report.py YYYY-MM-DD
(end date = last day of 7-day window, inclusive)
"""

import os
import sys
import sqlite3
from io import BytesIO
from datetime import datetime, timedelta
from statistics import mean, pstdev

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
DATA_DIR = os.environ.get("DATA_DIR", ROOT)
DB_PATH = os.path.join(DATA_DIR, "db", "clio.db")
TEMPLATE_PATH = os.path.join(ROOT, "templates", "base_template.pptx")
REPORTS_ROOT = os.path.join(DATA_DIR, "reports")

BG = RGBColor(0x0F, 0x17, 0x2A)
TEXT = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
AMBER = RGBColor(0xFB, 0x92, 0x3C)
RED = RGBColor(0xF8, 0x71, 0x71)
PANEL = RGBColor(0x1E, 0x29, 0x3B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def hex_to_rgb(hx):
    hx = (hx or "#3ECFCF").lstrip("#")
    if len(hx) != 6:
        return RGBColor(0x3E, 0xCF, 0xCF)
    return RGBColor(int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))


def fetch(end_date_str):
    end = datetime.strptime(end_date_str, "%Y-%m-%d").date()
    days = [(end - timedelta(days=6 - i)).strftime("%Y-%m-%d") for i in range(7)]

    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cur = conn.cursor()
    cur.execute("SELECT * FROM departments ORDER BY id")
    depts = [dict(r) for r in cur.fetchall()]

    by_dept = {d["id"]: {day: None for day in days} for d in depts}
    highlights = {d["id"]: [] for d in depts}

    cur.execute(
        """SELECT * FROM daily_submissions
           WHERE submission_date BETWEEN ? AND ? AND is_submitted = 1""",
        (days[0], days[-1]),
    )
    for r in cur.fetchall():
        row = dict(r)
        if row["department_id"] in by_dept and row["submission_date"] in by_dept[row["department_id"]]:
            by_dept[row["department_id"]][row["submission_date"]] = int(row["overall_progress"] or 0)
            if row.get("highlights"):
                highlights[row["department_id"]].append((row["submission_date"], row["highlights"]))
    conn.close()
    return depts, days, by_dept, highlights


def new_presentation():
    if os.path.exists(TEMPLATE_PATH):
        try:
            return Presentation(TEMPLATE_PATH)
        except Exception:
            pass
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return slide


def add_text(slide, left, top, width, height, text, *, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text if text is not None else ""
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb


def add_rect(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def cell_color_for(pct):
    if pct is None:
        return RGBColor(0x33, 0x33, 0x33)
    if pct >= 80:
        return GREEN
    if pct >= 50:
        return AMBER
    return RED


def slide_summary(prs, days, depts, by_dept):
    slide = blank_slide(prs)
    event = os.environ.get("EVENT_NAME", "Event")
    add_text(slide, Inches(0.4), Inches(0.4), Inches(12.5), Inches(0.7),
             f"Weekly Progress Report: {event}", size=26, bold=True)
    add_text(slide, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.4),
             f"Date range: {days[0]} → {days[-1]}", size=13, color=MUTED)

    # Table header
    x0 = Inches(0.4)
    y0 = Inches(1.8)
    col_dept_w = Inches(2.2)
    col_day_w = Inches(1.25)
    col_avg_w = Inches(1.3)
    header_h = Inches(0.4)
    row_h = Inches(0.4)

    add_rect(slide, x0, y0, Inches(13), header_h, PANEL)
    add_text(slide, x0 + Inches(0.1), y0 + Inches(0.08), col_dept_w, Inches(0.3),
             "DEPARTMENT", size=10, bold=True, color=MUTED)
    for i, day in enumerate(days):
        add_text(slide, x0 + col_dept_w + col_day_w * i + Inches(0.05),
                 y0 + Inches(0.08), col_day_w, Inches(0.3),
                 day[5:], size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, x0 + col_dept_w + col_day_w * 7 + Inches(0.05),
             y0 + Inches(0.08), col_avg_w, Inches(0.3),
             "AVG", size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    for r_idx, dept in enumerate(depts):
        top = y0 + header_h + row_h * r_idx
        if r_idx % 2 == 0:
            add_rect(slide, x0, top, Inches(13), row_h, RGBColor(0x0A, 0x0C, 0x12))
        add_rect(slide, x0, top, Inches(0.08), row_h, hex_to_rgb(dept["stream_color"]))
        add_text(slide, x0 + Inches(0.2), top + Inches(0.08),
                 col_dept_w, Inches(0.3), dept["name"], size=11, bold=True, color=TEXT)
        vals = []
        for i, day in enumerate(days):
            pct = by_dept[dept["id"]][day]
            cell_x = x0 + col_dept_w + col_day_w * i + Inches(0.1)
            cell_color = cell_color_for(pct)
            if pct is not None:
                vals.append(pct)
                add_text(slide, cell_x, top + Inches(0.08), col_day_w - Inches(0.2), Inches(0.3),
                         f"{pct}%", size=11, color=cell_color, align=PP_ALIGN.CENTER)
            else:
                add_text(slide, cell_x, top + Inches(0.08), col_day_w - Inches(0.2), Inches(0.3),
                         "n/a", size=11, color=MUTED, align=PP_ALIGN.CENTER)
        avg = round(mean(vals)) if vals else 0
        add_text(slide, x0 + col_dept_w + col_day_w * 7 + Inches(0.05),
                 top + Inches(0.08), col_avg_w, Inches(0.3),
                 f"{avg}%", size=12, bold=True, color=TEXT, align=PP_ALIGN.CENTER)


def make_trend_png(days, values, color_hex):
    fig, ax = plt.subplots(figsize=(9, 3.8), dpi=150)
    fig.patch.set_facecolor("#06080D")
    ax.set_facecolor("#06080D")
    ys = [v if v is not None else 0 for v in values]
    ax.plot(days, ys, marker="o", color=color_hex, linewidth=2.5, markersize=8)
    for i, v in enumerate(values):
        if v is None:
            ax.annotate("No Data", (days[i], 0), color="#6B7280",
                        ha="center", va="bottom", fontsize=8)
    ax.set_ylim(0, 105)
    ax.set_ylabel("Completion %", color="#6B7280")
    ax.tick_params(colors="#6B7280")
    for spine in ax.spines.values():
        spine.set_color("#1f2937")
    ax.grid(color="#1f2937", linestyle="--", linewidth=0.5)
    plt.xticks(rotation=30, ha="right")
    plt.tight_layout()
    buf = BytesIO()
    plt.savefig(buf, format="png", facecolor=fig.get_facecolor())
    plt.close(fig)
    buf.seek(0)
    return buf


def slide_dept_trend(prs, dept, days, by_dept, highlights):
    slide = blank_slide(prs)
    color = hex_to_rgb(dept["stream_color"])
    add_rect(slide, 0, 0, Inches(0.1), SLIDE_H, color)
    add_text(slide, Inches(0.4), Inches(0.4), Inches(12), Inches(0.5),
             f"{dept['name'].upper()}  •  WEEKLY TREND", size=22, bold=True, color=color)

    values = [by_dept[dept["id"]][d] for d in days]
    try:
        buf = make_trend_png(days, values, dept["stream_color"])
        slide.shapes.add_picture(buf, Inches(0.4), Inches(1.2),
                                 width=Inches(9.5), height=Inches(3.8))
    except Exception as e:
        sys.stderr.write(f"trend err: {e}\n")

    # Highlights list
    add_text(slide, Inches(0.4), Inches(5.2), Inches(12), Inches(0.3),
             "HIGHLIGHTS THIS WEEK", size=11, bold=True, color=MUTED)
    hlist = highlights.get(dept["id"], [])
    if hlist:
        text = "\n".join([f"• {d}: {h}" for d, h in hlist[:6]])
    else:
        text = "No highlights logged this week."
    add_text(slide, Inches(0.4), Inches(5.55), Inches(12.5), Inches(1.8),
             text, size=11, color=TEXT)


def slide_insights(prs, depts, days, by_dept):
    slide = blank_slide(prs)
    add_text(slide, Inches(0.4), Inches(0.4), Inches(12), Inches(0.6),
             "WEEKLY INSIGHTS", size=26, bold=True)

    stats = []
    for d in depts:
        vals = [by_dept[d["id"]][day] for day in days if by_dept[d["id"]][day] is not None]
        missing = sum(1 for day in days if by_dept[d["id"]][day] is None)
        if vals:
            stats.append({
                "name": d["name"], "avg": mean(vals),
                "stdev": pstdev(vals) if len(vals) > 1 else 0,
                "missing": missing,
            })
        else:
            stats.append({"name": d["name"], "avg": 0, "stdev": 999, "missing": 7})

    most_consistent = min(stats, key=lambda s: s["stdev"]) if stats else None
    highest = max(stats, key=lambda s: s["avg"]) if stats else None
    missing_some = [s for s in stats if s["missing"] > 0]

    y = Inches(1.4)
    if most_consistent:
        add_text(slide, Inches(0.4), y, Inches(12), Inches(0.4),
                 f"Most consistent department: {most_consistent['name']} "
                 f"(σ = {most_consistent['stdev']:.1f})",
                 size=15, color=GREEN)
        y += Inches(0.5)
    if highest:
        add_text(slide, Inches(0.4), y, Inches(12), Inches(0.4),
                 f"Highest average completion: {highest['name']} ({highest['avg']:.0f}%)",
                 size=15, color=GREEN)
        y += Inches(0.5)

    y += Inches(0.3)
    add_text(slide, Inches(0.4), y, Inches(12), Inches(0.4),
             "Departments with missing submissions this week:",
             size=13, bold=True, color=MUTED)
    y += Inches(0.4)
    if missing_some:
        for s in missing_some:
            add_text(slide, Inches(0.6), y, Inches(12), Inches(0.35),
                     f"• {s['name']}: {s['missing']} day(s) missing", size=12, color=RED)
            y += Inches(0.35)
    else:
        add_text(slide, Inches(0.6), y, Inches(12), Inches(0.35),
                 "None, full attendance this week. ✓", size=12, color=GREEN)


def build(end_date):
    depts, days, by_dept, highlights = fetch(end_date)
    prs = new_presentation()

    slide_summary(prs, days, depts, by_dept)
    for d in depts:
        slide_dept_trend(prs, d, days, by_dept, highlights)
    slide_insights(prs, depts, days, by_dept)

    total = len(prs.slides)
    for idx, slide in enumerate(prs.slides, start=1):
        add_text(slide, Inches(12.6), Inches(7.1), Inches(0.7), Inches(0.3),
                 f"{idx} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)

    out_dir = os.path.join(REPORTS_ROOT, "weekly")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, f"weekly_report_{end_date}.pptx")
    prs.save(out_path)
    return out_path


def main():
    if len(sys.argv) < 2:
        sys.stderr.write("usage: generate_weekly_report.py YYYY-MM-DD\n")
        sys.exit(1)
    end_date = sys.argv[1]
    try:
        datetime.strptime(end_date, "%Y-%m-%d")
    except ValueError:
        sys.stderr.write(f"invalid date: {end_date}\n")
        sys.exit(1)
    try:
        print(build(end_date))
    except Exception as e:
        import traceback
        traceback.print_exc(file=sys.stderr)
        sys.stderr.write(f"weekly generation failed: {e}\n")
        sys.exit(1)


if __name__ == "__main__":
    main()
